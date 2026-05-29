"""Generate docs/agentic-architecture-deck.pptx — the agentic-ops architecture
deck for the Predictive RPC Estimator.

Audience: technical + product stakeholders in a global corporate setting.
Covers the new reasoning-agent architecture (ADR 0005) and how to run the demo.

Design (per brand): plain white slide backgrounds; flat shapes, icons and
fills — no 3D, gradients or lighting; Poppins throughout; colours drawn only
from the brand palette led by #0064FF.

Idempotent: re-run to regenerate.
"""
from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT = Path(__file__).parent / "agentic-architecture-deck.pptx"
PROJECT = "msm-rpc"
REGION = "europe-west2"
BREAKER_URL = "https://breaker-agent-staging-ifcjcfl7xa-nw.a.run.app"
BOUNDS_URL = "https://bounds-agent-staging-ifcjcfl7xa-nw.a.run.app"
DRIFT_URL = "https://drift-agent-staging-ifcjcfl7xa-nw.a.run.app"

# Palette — flat, brand-blue led. Every colour is a member of the brand list.
NAVY = RGBColor(0x00, 0x26, 0x59)   # #002659 deepest brand blue
DEEP = RGBColor(0x05, 0x16, 0x2E)   # #05162e near-black blue
TEAL = RGBColor(0x00, 0x64, 0xFF)   # #0064FF primary brand
ACCENT = RGBColor(0x00, 0x64, 0xFF)  # brand (accent == brand)
BLUE2 = RGBColor(0x42, 0x85, 0xF4)  # #4285f4 secondary blue
STEEL = RGBColor(0x45, 0x60, 0x85)  # #456085 muted brand blue
SKY = RGBColor(0x79, 0xAD, 0xFF)    # #79adff light brand blue
SLATE = RGBColor(0x21, 0x21, 0x21)  # #212121 body text
MUTED = RGBColor(0x7D, 0x7D, 0x7D)  # #7d7d7d secondary text
LIGHT = RGBColor(0xE7, 0xF0, 0xFE)  # #e7f0fe card tint
TINT = RGBColor(0xE2, 0xEE, 0xFF)   # #e2eeff lighter tint
RULE = RGBColor(0xD2, 0xE3, 0xFC)   # #d2e3fc hairline
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Poppins"
TOTAL_PAGES = 10


def run(p, text, size=14, bold=False, color=SLATE, font=FONT, italic=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color
    return r


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def title_bar(slide, prs, title, eyebrow=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    fill(bar, NAVY)
    tf = bar.text_frame
    tf.margin_left = Inches(0.6); tf.margin_top = Inches(0.16)
    tf.word_wrap = True
    if eyebrow:
        p = tf.paragraphs[0]
        run(p, eyebrow.upper(), size=10, bold=True, color=SKY)
        p2 = tf.add_paragraph()
        run(p2, title, size=24, bold=True, color=WHITE)
    else:
        p = tf.paragraphs[0]
        run(p, title, size=26, bold=True, color=WHITE)
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.95),
                                 Inches(1.2), Pt(3))
    fill(acc, ACCENT)


def footer(slide, prs, page_num, label="Predictive RPC Estimator — Agentic Ops"):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.35),
        prs.slide_width, Inches(0.35),
    )
    fill(bar, LIGHT)
    tf = bar.text_frame
    tf.margin_left = Inches(0.6); tf.margin_top = Inches(0.06)
    p = tf.paragraphs[0]
    run(p, label, size=9, color=MUTED)
    pn = slide.shapes.add_textbox(
        prs.slide_width - Inches(1.2), prs.slide_height - Inches(0.34),
        Inches(1.0), Inches(0.3),
    )
    pp = pn.text_frame.paragraphs[0]
    pp.alignment = PP_ALIGN.RIGHT
    run(pp, f"{page_num} / {TOTAL_PAGES}", size=9, color=MUTED)


def hero(slide, prs, big, sub=None, eyebrow=None):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.5),
                                  prs.slide_width, Inches(0.05))
    fill(band, TEAL)
    if eyebrow:
        eb = slide.shapes.add_textbox(Inches(0.7), Inches(1.5),
                                      prs.slide_width - Inches(1.4), Inches(0.4))
        run(eb.text_frame.paragraphs[0], eyebrow.upper(), size=13, bold=True, color=TEAL)
    t = slide.shapes.add_textbox(Inches(0.7), Inches(2.0),
                                 prs.slide_width - Inches(1.4), Inches(2.3))
    tf = t.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.line_spacing = 1.05
    run(p, big, size=50, bold=True, color=NAVY)
    if sub:
        st = slide.shapes.add_textbox(Inches(0.7), Inches(4.85),
                                      prs.slide_width - Inches(1.4), Inches(2.0))
        tf = st.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.line_spacing = 1.35
        run(p, sub, size=21, color=SLATE)


def card(slide, x, y, w, h, tint=WHITE, border=True):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    c.adjustments[0] = 0.05
    fill(c, tint)
    if border:
        c.line.color.rgb = RULE; c.line.width = Pt(0.75)
    return c


def code_card(slide, x, y, w, h, lines):
    """A terminal-style card. Poppins throughout (per brand), tinted dark blue."""
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    c.adjustments[0] = 0.05
    fill(c, DEEP)
    tf = c.text_frame
    tf.margin_left = Inches(0.22); tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.14); tf.margin_bottom = Inches(0.10)
    tf.word_wrap = True
    for i, (kind, text) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.18
        if kind == "cmd":
            run(p, "$ ", size=11, bold=True, color=SKY)
            run(p, text, size=11, color=WHITE)
        elif kind == "out":
            run(p, text, size=10, color=SKY)
        elif kind == "cmt":
            run(p, text, size=10, italic=True, color=RULE)
    return c


# ───────────────────────── Slides ─────────────────────────

def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.5),
                              prs.slide_width, Inches(0.05))
    fill(band, TEAL)
    eb = s.shapes.add_textbox(Inches(0.7), Inches(1.55), Inches(12), Inches(0.4))
    run(eb.text_frame.paragraphs[0],
        "AGENTIC OPS — REASONING AGENTS ON THE PREDICTIVE RPC ESTIMATOR",
        size=12, bold=True, color=TEAL)
    t = s.shapes.add_textbox(Inches(0.7), Inches(2.05), Inches(12), Inches(1.6))
    tf = t.text_frame; tf.word_wrap = True
    run(tf.paragraphs[0], "From fixed guardrails to judgment", size=44, bold=True, color=NAVY)
    p2 = tf.add_paragraph()
    run(p2, "Three ADK reasoning agents on the live scoring platform", size=20, color=SLATE)
    sub = s.shapes.add_textbox(Inches(0.7), Inches(4.85), Inches(12), Inches(2))
    tf = sub.text_frame; tf.word_wrap = True
    for label, val in [
        ("Scope", "Breaker triage · Bounds calibration · Drift triage"),
        ("Runtime", "Google ADK on Cloud Run · europe-west2 · Gemini"),
        ("Status", "Live on staging · reasoning verified against the model · CI green"),
    ]:
        p = tf.add_paragraph(); p.line_spacing = 1.3
        run(p, f"{label}    ", size=12, bold=True, color=TEAL)
        run(p, val, size=14, color=SLATE)


def slide_shift(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    hero(
        s, prs,
        big="The platform already decides. Now it can reason.",
        sub=("The deterministic guardrails — bounds, breaker, drift scoring — are fast and "
             "provably correct, but they can't judge. Three places needed a human: triage an "
             "incident, trust a bounds change, act on drift. That judgment is now an agent."),
        eyebrow="What changed",
    )
    bottom = s.shapes.add_textbox(Inches(0.7), Inches(6.15),
                                  prs.slide_width - Inches(1.4), Inches(0.7))
    tf = bottom.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.line_spacing = 1.3
    for i, (k, v) in enumerate([
        ("UNCHANGED", "Scoring hot path, reconciliation, activation — no LLM, no latency tax."),
        ("ADDED", "Reasoning only where judgment lived."),
        ("SAFE", "Every agent write is schema-gated before it touches state."),
    ]):
        if i > 0:
            run(p, "      ", size=11)
        run(p, k + "  ", size=10, bold=True, color=TEAL)
        run(p, v, size=11, color=SLATE)
    footer(s, prs, 2)


def slide_principle(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, prs, "The LLM reasons. The domain decides.",
              eyebrow="The one rule that makes this safe")
    cols = [
        ("Deterministic core stays", STEEL,
         "should_trip, propose_bounds, drift PSI verdicts remain pure, tested functions. "
         "The agent never recomputes them — it calls them as read-only tools."),
        ("The LLM interprets", TEAL,
         "It reads the tool outputs and decides what they mean and which action to take — "
         "the judgment a human used to supply, now encoded and observable."),
        ("The write is schema-gated", NAVY,
         "Agent output is validated against an explicit schema, then a domain value object, "
         "before any retrain / PR / page. The model cannot act on free text."),
    ]
    cw = (prs.slide_width - Inches(1.1) - Inches(0.5)) / 3
    ch = Inches(4.0); top = Inches(1.45); left0 = Inches(0.55); gap = Inches(0.25)
    for i, (head, color, body) in enumerate(cols):
        x = left0 + i * (cw + gap)
        card(s, x, top, cw, ch)
        stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top, cw, Inches(0.08))
        fill(stripe, color)
        num = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), top + Inches(0.35),
                                 Inches(0.6), Inches(0.6))
        fill(num, color)
        ntf = num.text_frame; ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
        np = ntf.paragraphs[0]; np.alignment = PP_ALIGN.CENTER
        run(np, str(i + 1), size=22, bold=True, color=WHITE)
        tb = s.shapes.add_textbox(x + Inches(0.3), top + Inches(1.15),
                                  cw - Inches(0.6), ch - Inches(1.3))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.line_spacing = 1.1
        run(p, head, size=18, bold=True, color=NAVY)
        p2 = tf.add_paragraph(); p2.space_before = Pt(10); p2.line_spacing = 1.32
        run(p2, body, size=12, color=SLATE)
    note = s.shapes.add_textbox(Inches(0.55), Inches(5.75),
                                prs.slide_width - Inches(1.1), Inches(0.6))
    p = note.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run(p, "Result: domain logic keeps its ≥95% test floor; nondeterminism is contained to the "
        "reasoning layer and verified by evals.", size=11, italic=True, color=MUTED)
    footer(s, prs, 3)


def slide_three_agents(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, prs, "Three agents, one bounded context each",
              eyebrow="What each reasons about — and what stays deterministic")
    agents = [
        ("Breaker triage", TEAL,
         "Decides escalation after a trip: severity, page vs annotate.",
         "should_trip → engage (trip path untouched)",
         "Pub/Sub push on rpc-anomaly"),
        ("Bounds calibration", BLUE2,
         "Judges a proposed bounds change: genuine shift vs transient vs noise.",
         "propose_bounds (the numbers)",
         "Cloud Scheduler · weekly"),
        ("Drift triage", NAVY,
         "Decides retrain / alert / noop and names the driver features.",
         "detect_drift + PSI verdicts",
         "Cloud Scheduler · daily"),
    ]
    cw = (prs.slide_width - Inches(1.1) - Inches(0.5)) / 3
    ch = Inches(4.45); top = Inches(1.30); left0 = Inches(0.55); gap = Inches(0.25)
    for i, (name, color, reasons, deterministic, trigger) in enumerate(agents):
        x = left0 + i * (cw + gap)
        card(s, x, top, cw, ch)
        head = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top, cw, Inches(0.6))
        fill(head, color)
        htf = head.text_frame; htf.margin_left = Inches(0.2); htf.vertical_anchor = MSO_ANCHOR.MIDDLE
        run(htf.paragraphs[0], name, size=16, bold=True, color=WHITE)
        tb = s.shapes.add_textbox(x + Inches(0.22), top + Inches(0.8),
                                  cw - Inches(0.44), ch - Inches(0.95))
        tf = tb.text_frame; tf.word_wrap = True
        run(tf.paragraphs[0], "REASONS ABOUT", size=9, bold=True, color=color)
        p = tf.add_paragraph(); p.line_spacing = 1.28; p.space_after = Pt(10)
        run(p, reasons, size=12.5, color=SLATE)
        p = tf.add_paragraph()
        run(p, "STAYS DETERMINISTIC", size=9, bold=True, color=color)
        p = tf.add_paragraph(); p.line_spacing = 1.25; p.space_after = Pt(10)
        run(p, deterministic, size=11, color=SLATE)
        p = tf.add_paragraph()
        run(p, "TRIGGER", size=9, bold=True, color=color)
        p = tf.add_paragraph(); p.line_spacing = 1.25
        run(p, trigger, size=11, color=SLATE)
    footer(s, prs, 4)


def slide_architecture(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, prs, "How an agent is wired",
              eyebrow="Same shape for all three — trigger to gated write")
    lanes = [
        ("Trigger", STEEL, [
            ("Pub/Sub push", "breaker — rpc-anomaly topic"),
            ("Cloud Scheduler", "bounds (weekly) · drift (daily)"),
            ("FastAPI on Cloud Run", "scales to zero · /health"),
        ]),
        ("Agent (ADK)", TEAL, [
            ("Investigator", "read-only tools gather evidence"),
            ("Decider", "forced output schema, no tools"),
            ("Gemini", "via API key — region-independent"),
        ]),
        ("Read-only tools", BLUE2, [
            ("Deterministic core", "should_trip · propose_bounds · drift PSI"),
            ("BigQuery", "psi_daily · rpc_predictions"),
            ("Model registry", "latest model version"),
        ]),
        ("Gated write", NAVY, [
            ("Schema gate", "Pydantic → domain value object (§4)"),
            ("Dispatch", "page · open PR · trigger retrain Job"),
            ("Per-AI-call log", "model · tokens · cost (§6)"),
        ]),
    ]
    n = len(lanes)
    margin = Inches(0.35)
    avail_w = prs.slide_width - Inches(0.8)
    lane_w = (avail_w - margin * (n - 1)) / n
    top = Inches(1.25); lane_h = Inches(5.35)
    for i, (lane_title, accent, items) in enumerate(lanes):
        x = Inches(0.4) + i * (lane_w + margin)
        card(s, x, top, lane_w, lane_h)
        stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top, lane_w, Inches(0.5))
        fill(stripe, accent)
        hb = s.shapes.add_textbox(x + Inches(0.15), top + Inches(0.06),
                                  lane_w - Inches(0.3), Inches(0.4))
        run(hb.text_frame.paragraphs[0], lane_title.upper(), size=11, bold=True, color=WHITE)
        item_top = top + Inches(0.65); item_h = Inches(1.35); item_gap = Inches(0.16)
        for j, (label, sub) in enumerate(items):
            iy = item_top + j * (item_h + item_gap)
            ic = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    x + Inches(0.15), iy, lane_w - Inches(0.3), item_h)
            ic.adjustments[0] = 0.10
            fill(ic, LIGHT); ic.line.color.rgb = RULE; ic.line.width = Pt(0.5)
            tf = ic.text_frame
            tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.10)
            tf.margin_top = Inches(0.14); tf.margin_bottom = Inches(0.08)
            tf.word_wrap = True
            run(tf.paragraphs[0], label, size=11, bold=True, color=NAVY)
            p2 = tf.add_paragraph(); p2.line_spacing = 1.15
            run(p2, sub, size=9, color=SLATE)
        if i < n - 1:
            # Integer EMU only — fractional coords from float division make
            # add_connector emit invalid OOXML (PowerPoint flags it corrupt).
            arrow_y = int(top + lane_h / 2)
            ln = s.shapes.add_connector(2, int(x + lane_w + Pt(2)), arrow_y,
                                        int(x + lane_w + margin - Pt(2)), arrow_y)
            ln.line.color.rgb = MUTED; ln.line.width = Pt(2)
    footer(s, prs, 5)


def slide_governance(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, prs, "Safe by construction",
              eyebrow="Why a non-deterministic agent is allowed to act")
    blocks = [
        ("Writes are schema-gated", TEAL, [
            "Agent emits a validated structured output, mapped to a domain value object",
            "Invariants enforced in the constructor — bad output raises before any side effect",
            "The LLM cannot retrain, open a PR or page on free text (§4)",
        ]),
        ("Least-privilege per agent", BLUE2, [
            "One service account per agent; scoped IAM in Terraform",
            "Breaker can't trip the breaker; bounds can't retrain; none hold broad keys",
            "Self-hosted on Cloud Run, scale-to-zero — cost only when triggered",
        ]),
        ("Verified reasoning", NAVY, [
            "Hermetic eval suites run each agent against the real model per scenario",
            "Asserts discrimination: escalate vs not, retrain vs not, PR vs not",
            "Runs in CI; transient model flakiness absorbed by reruns (ADR 0006)",
        ]),
        ("Observable + reversible", STEEL, [
            "Per-AI-call log: model, prompt hash, tokens, cost (§6)",
            "Deterministic cores keep their ≥95% test floor — unchanged",
            "Gated behind DEPLOY_AGENTS; off = removed, re-tag = redeployed",
        ]),
    ]
    cw = (prs.slide_width - Inches(1.1) - Inches(0.25)) / 2
    ch = Inches(2.6); left0 = Inches(0.55); top0 = Inches(1.20); gap = Inches(0.25)
    for i, (label, color, items) in enumerate(blocks):
        row, col = i // 2, i % 2
        x = left0 + col * (cw + gap); y = top0 + row * (ch + gap)
        card(s, x, y, cw, ch)
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.25), y + Inches(0.25),
                                   Inches(0.55), Inches(0.55))
        fill(badge, color)
        btf = badge.text_frame; btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
        run(bp, "✓", size=20, bold=True, color=WHITE)
        hb = s.shapes.add_textbox(x + Inches(0.95), y + Inches(0.30),
                                  cw - Inches(1.15), Inches(0.5))
        run(hb.text_frame.paragraphs[0], label, size=16, bold=True, color=NAVY)
        ib = s.shapes.add_textbox(x + Inches(0.25), y + Inches(0.95),
                                  cw - Inches(0.50), ch - Inches(1.05))
        tf = ib.text_frame; tf.word_wrap = True
        for j, it in enumerate(items):
            pp = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            pp.line_spacing = 1.22; pp.space_after = Pt(3)
            run(pp, "•  ", size=10, bold=True, color=color)
            run(pp, it, size=10, color=SLATE)
    footer(s, prs, 6)


def slide_deploy(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, prs, "Live on staging", eyebrow="Cloud Run · europe-west2 · opt-in gated")
    rows = [
        ("Breaker triage", BREAKER_URL, "Pub/Sub push · rpc-anomaly", TEAL),
        ("Bounds calibration", BOUNDS_URL, "Cloud Scheduler · Mon 06:00 UTC", BLUE2),
        ("Drift triage", DRIFT_URL, "Cloud Scheduler · daily 06:30 UTC", NAVY),
    ]
    top0 = Inches(1.35); row_h = Inches(1.15); gap = Inches(0.20)
    for i, (name, url, trig, color) in enumerate(rows):
        y = top0 + i * (row_h + gap)
        card(s, Inches(0.55), y, prs.slide_width - Inches(1.1), row_h)
        stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), y, Inches(0.10), row_h)
        fill(stripe, color)
        tb = s.shapes.add_textbox(Inches(0.85), y + Inches(0.14),
                                  prs.slide_width - Inches(1.6), row_h - Inches(0.28))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        run(p, name + "    ", size=16, bold=True, color=NAVY)
        run(p, trig, size=11, color=MUTED)
        p2 = tf.add_paragraph(); p2.space_before = Pt(4)
        run(p2, url, size=12, color=TEAL)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(5.65),
                             prs.slide_width - Inches(1.1), Inches(0.95))
    fill(bar, DEEP)
    tf = bar.text_frame; tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.14)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run(p, "OPT-IN, ZERO COST UNTIL ON    ", size=10, bold=True, color=SKY)
    run(p, "DEPLOY_AGENTS=true + a version tag builds the images and applies the Terraform.",
        size=11, color=WHITE)
    p2 = tf.add_paragraph(); p2.space_before = Pt(6)
    run(p2, "Each service scales to zero and is auth-locked — only its own SA, the scheduler "
        "and Pub/Sub can invoke it.", size=11, color=LIGHT)
    footer(s, prs, 7)


def slide_run_demo(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, prs, "How to run the demo", eyebrow="Five steps · one terminal")
    # Left: numbered beats. Right: the commands.
    beats = [
        ("1", "Show it's live", "Three Cloud Run services, Ready, scaled to zero."),
        ("2", "Breaker triage, end to end",
         "Publish a synthetic anomaly; the agent reasons and escalates to the incident topic."),
        ("3", "Drift triage", "Fire the daily job; watch the retrain/alert/noop decision in the logs."),
        ("4", "Bounds calibration", "Fire the weekly job; watch the genuine-shift vs transient verdict."),
        ("5", "Reasoning quality", "Run the hermetic eval suites against the live model."),
    ]
    bx = Inches(0.55); bw = Inches(5.0); top = Inches(1.30)
    bh = Inches(1.02); gap = Inches(0.12)
    colors = [TEAL, TEAL, BLUE2, BLUE2, NAVY]
    for i, (num, head, sub) in enumerate(beats):
        y = top + i * (bh + gap)
        card(s, bx, y, bw, bh)
        nb = s.shapes.add_shape(MSO_SHAPE.OVAL, bx + Inches(0.18), y + Inches(0.26),
                                Inches(0.5), Inches(0.5))
        fill(nb, colors[i])
        ntf = nb.text_frame; ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
        npar = ntf.paragraphs[0]; npar.alignment = PP_ALIGN.CENTER
        run(npar, num, size=18, bold=True, color=WHITE)
        tb = s.shapes.add_textbox(bx + Inches(0.85), y + Inches(0.12), bw - Inches(1.0), bh - Inches(0.2))
        tf = tb.text_frame; tf.word_wrap = True
        run(tf.paragraphs[0], head, size=14, bold=True, color=NAVY)
        p2 = tf.add_paragraph(); p2.line_spacing = 1.18
        run(p2, sub, size=10.5, color=SLATE)
    # Right: commands
    cx = Inches(5.75); cw = prs.slide_width - cx - Inches(0.55)
    code_card(s, cx, top, cw, Inches(5.55), [
        ("cmt", "# 1 — it's live"),
        ("cmd", "gcloud run services list --project msm-rpc \\"),
        ("cmd", "  --region europe-west2 --filter name~agent"),
        ("cmt", "# 2 — breaker triage end-to-end"),
        ("cmd", "gcloud pubsub topics publish rpc-anomaly-staging \\"),
        ("cmd", "  --project msm-rpc --message \\"),
        ("cmd", "  '{\"kind\":\"null_or_zero_rate\",\"value\":0.08,"),
        ("cmd", "   \"threshold\":0.03,\"occurred_at_ms\":1717000000000}'"),
        ("cmt", "# 3 — drift   # 4 — bounds"),
        ("cmd", "gcloud scheduler jobs run drift-agent-staging \\"),
        ("cmd", "  --project msm-rpc --location europe-west2"),
        ("cmt", "# watch any agent's reasoning"),
        ("cmd", "gcloud logging read \\"),
        ("cmd", "  'resource.labels.service_name=\"breaker-agent-staging\"' \\"),
        ("cmd", "  --project msm-rpc --limit 10 --freshness 5m"),
        ("cmt", "# 5 — reasoning verified against the model"),
        ("cmd", "pytest -k eval -v   # in each agent service"),
    ])
    footer(s, prs, 8)


def slide_proof(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, prs, "Proven, not just plausible",
              eyebrow="Each agent discriminates the right way — against the real model")
    headers = ["Agent", "High-action scenario", "Low-action scenario"]
    rows = [
        ("Breaker triage", "Severe sustained breach → escalate", "Marginal blip → not a SEV1 page"),
        ("Bounds calibration", "Sustained shift → opens a PR", "Single spike day → declines"),
        ("Drift triage", "Multi-feature breach → retrain / alert", "Single warn → no retrain"),
    ]
    tx = Inches(0.55); tw = prs.slide_width - Inches(1.1); top = Inches(1.40)
    colw = [Inches(3.2), tw / 2 - Inches(1.0), tw / 2 - Inches(1.0)]
    # header row
    hh = Inches(0.5)
    hbar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, tx, top, tw, hh)
    fill(hbar, NAVY)
    cx = tx
    for k, h in enumerate(headers):
        hb = s.shapes.add_textbox(cx + Inches(0.15), top + Inches(0.06), colw[k], Inches(0.4))
        run(hb.text_frame.paragraphs[0], h.upper(), size=11, bold=True, color=WHITE)
        cx += colw[k]
    # body rows
    rh = Inches(0.95)
    for i, r in enumerate(rows):
        y = top + hh + i * rh
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, tx, y, tw, rh)
        fill(bg, WHITE if i % 2 == 0 else TINT)
        bg.line.color.rgb = RULE; bg.line.width = Pt(0.5)
        cx = tx
        for k, cell in enumerate(r):
            cb = s.shapes.add_textbox(cx + Inches(0.15), y + Inches(0.18), colw[k] - Inches(0.2), rh - Inches(0.3))
            tf = cb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.line_spacing = 1.15
            if k == 0:
                run(p, cell, size=13, bold=True, color=NAVY)
            else:
                run(p, "✓  ", size=12, bold=True, color=TEAL)
                run(p, cell, size=12, color=SLATE)
            cx += colw[k]
    note = s.shapes.add_textbox(Inches(0.55), Inches(5.55), tw, Inches(1.0))
    tf = note.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.line_spacing = 1.3
    run(p, "Verified in CI against Gemini — the eval job is green on every push. ",
        size=12, color=SLATE)
    run(p, "Deterministic cores carry their own unit tests at the ≥95% / ≥85% floors.",
        size=12, bold=True, color=TEAL)
    footer(s, prs, 9)


def slide_next(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.5),
                              prs.slide_width, Inches(0.05))
    fill(band, TEAL)
    run(s.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12), Inches(0.4))
        .text_frame.paragraphs[0], "WHAT'S NEXT", size=13, bold=True, color=TEAL)
    t = s.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(12), Inches(1.3))
    run(t.text_frame.paragraphs[0], "Live, gated, and ready for signal.", size=46, bold=True, color=NAVY)
    sub = s.shapes.add_textbox(Inches(0.7), Inches(3.85), Inches(12), Inches(2.6))
    tf = sub.text_frame; tf.word_wrap = True
    for k, v in [
        ("Add the GitHub token", "so the bounds agent can open recalibration PRs (Contents + PR write)."),
        ("Materialize psi_daily", "run Dataform in staging — until then drift & signal correlation see no data."),
        ("Then watch them work", "real anomalies, real drift, real bounds shifts drive the agents end-to-end."),
    ]:
        p = tf.add_paragraph(); p.line_spacing = 1.4
        run(p, "•  ", size=16, bold=True, color=TEAL)
        run(p, k + " — ", size=16, bold=True, color=NAVY)
        run(p, v, size=15, color=SLATE)
    foot = s.shapes.add_textbox(Inches(0.7), Inches(6.8), Inches(12), Inches(0.4))
    run(foot.text_frame.paragraphs[0],
        "Agents live on staging  ·  europe-west2  ·  reasoning verified  ·  CI green",
        size=11, color=MUTED)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_cover(prs)          # 1
    slide_shift(prs)          # 2
    slide_principle(prs)      # 3
    slide_three_agents(prs)   # 4
    slide_architecture(prs)   # 5
    slide_governance(prs)     # 6
    slide_deploy(prs)         # 7
    slide_run_demo(prs)       # 8
    slide_proof(prs)          # 9
    slide_next(prs)           # 10
    prs.save(OUT)
    print(f"wrote {OUT}  ({OUT.stat().st_size:,} bytes, {len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
