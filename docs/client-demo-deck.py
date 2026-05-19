"""Generate docs/client-demo-deck.pptx — client-facing demo deck for the
Predictive RPC Estimator.

Audience: client executive + technical stakeholders attending the
solution demo. Tone: outcome-led, confident, light on internal jargon.

Idempotent: re-run to regenerate.
"""
from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT = Path(__file__).parent / "client-demo-deck.pptx"
DEMO_URL = "https://dashboard-staging-794974391956.europe-west2.run.app"
API_URL  = "https://scoring-api-staging-794974391956.europe-west2.run.app"

# Palette — modern corporate; brand-blue led; flat, no gradients/3D.
# Variable names are kept from the previous deck for minimal churn;
# semantics shift: TEAL/ACCENT/GREEN now all map to brand-blue tones.
NAVY   = RGBColor(0x00, 0x26, 0x59)  # deepest brand blue — strong text + stripes
TEAL   = RGBColor(0x00, 0x64, 0xFF)  # primary brand
ACCENT = RGBColor(0x00, 0x64, 0xFF)  # primary brand (same — accent is the brand)
GREEN  = RGBColor(0x42, 0x85, 0xF4)  # secondary blue (was "success / complete")
SLATE  = RGBColor(0x21, 0x21, 0x21)  # body text — near-black
MUTED  = RGBColor(0x7D, 0x7D, 0x7D)  # secondary / muted text
LIGHT  = RGBColor(0xE7, 0xF0, 0xFE)  # card tint
RULE   = RGBColor(0xD2, 0xE3, 0xFC)  # border / hairline
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Poppins"
TOTAL_PAGES = 11


def run(p, text, size=14, bold=False, color=SLATE, font=FONT, italic=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color
    return r


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def line(slide, x1, y1, x2, y2, color=RULE, weight=0.75):
    ln = slide.shapes.add_connector(1, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def title_bar(slide, prs, title, eyebrow=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    fill(bar, NAVY)
    tf = bar.text_frame
    tf.margin_left = Inches(0.6); tf.margin_top = Inches(0.18)
    tf.word_wrap = True
    if eyebrow:
        p = tf.paragraphs[0]
        run(p, eyebrow.upper(), size=10, bold=True, color=ACCENT)
        p2 = tf.add_paragraph()
        run(p2, title, size=24, bold=True, color=WHITE)
    else:
        p = tf.paragraphs[0]
        run(p, title, size=26, bold=True, color=WHITE)
    # accent underline
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.95),
                                 Inches(1.2), Pt(3))
    fill(acc, ACCENT)


def footer(slide, prs, page_num, label="Predictive RPC Estimator — Client Demo"):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.35),
        prs.slide_width, Inches(0.35),
    )
    fill(bar, LIGHT)
    tf = bar.text_frame
    tf.margin_left = Inches(0.6); tf.margin_top = Inches(0.06)
    p = tf.paragraphs[0]
    run(p, f"{label}", size=9, color=MUTED)
    # right-aligned page number
    pn = slide.shapes.add_textbox(
        prs.slide_width - Inches(1.2), prs.slide_height - Inches(0.34),
        Inches(1.0), Inches(0.3),
    )
    pp = pn.text_frame.paragraphs[0]
    pp.alignment = PP_ALIGN.RIGHT
    run(pp, f"{page_num} / {TOTAL_PAGES}", size=9, color=MUTED)


def body_box(slide, prs, lines, top=Inches(1.25), left=Inches(0.6),
             width=None, height=None, font_size=14, line_spacing=1.2):
    width = width or prs.slide_width - Inches(1.2)
    height = height or prs.slide_height - Inches(1.7)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        if isinstance(ln, tuple):
            kind = ln[0]
            if kind == "h":
                run(p, ln[1], size=font_size + 4, bold=True, color=TEAL)
                p.space_after = Pt(6)
            elif kind == "b":
                run(p, "•  ", size=font_size, bold=True, color=ACCENT)
                run(p, ln[1], size=font_size, color=SLATE)
                p.space_after = Pt(3)
            elif kind == "sub":
                run(p, "    – ", size=font_size - 1, color=MUTED)
                run(p, ln[1], size=font_size - 1, color=SLATE)
            elif kind == "p":
                run(p, ln[1], size=font_size, color=SLATE)
                p.space_after = Pt(6)
            elif kind == "mute":
                run(p, ln[1], size=font_size - 2, color=MUTED, italic=True)
        else:
            run(p, ln, size=font_size, color=SLATE)
    return box


def metric_card(slide, left, top, w, h, value, label, value_color=NAVY, accent=TEAL):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    card.adjustments[0] = 0.06
    fill(card, WHITE)
    card.line.color.rgb = RULE
    card.line.width = Pt(0.75)
    # accent stripe
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, Inches(0.08))
    fill(stripe, accent)
    # value
    vb = slide.shapes.add_textbox(left, top + Inches(0.25), w, Inches(0.7))
    vp = vb.text_frame.paragraphs[0]; vp.alignment = PP_ALIGN.CENTER
    run(vp, value, size=28, bold=True, color=value_color)
    # label
    lb = slide.shapes.add_textbox(left, top + Inches(0.95), w, Inches(0.5))
    lp = lb.text_frame.paragraphs[0]; lp.alignment = PP_ALIGN.CENTER
    lp.line_spacing = 1.1
    run(lp, label, size=11, color=MUTED)

def hero_statement(slide, prs, big, sub=None, eyebrow=None,
                   bg_color=None, big_color=None, sub_color=None):
    """Full-bleed hero slide. One enormous statement, optional subline.
    Designed to land in the gut, not the head."""
    if bg_color is not None:
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    0, 0, prs.slide_width, prs.slide_height)
        fill(bg, bg_color)
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      0, Inches(4.20),
                                      prs.slide_width, Inches(0.06))
        fill(band, ACCENT)

    big_color = big_color or (WHITE if bg_color else NAVY)
    sub_color = sub_color or (LIGHT if bg_color else SLATE)

    if eyebrow:
        eb = slide.shapes.add_textbox(Inches(0.7), Inches(1.4),
                                      prs.slide_width - Inches(1.4), Inches(0.4))
        p = eb.text_frame.paragraphs[0]
        run(p, eyebrow.upper(), size=13, bold=True, color=ACCENT)

    t = slide.shapes.add_textbox(Inches(0.7), Inches(2.1),
                                 prs.slide_width - Inches(1.4), Inches(2.2))
    tf = t.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.line_spacing = 1.05
    run(p, big, size=54, bold=True, color=big_color)

    if sub:
        st = slide.shapes.add_textbox(Inches(0.7), Inches(4.7),
                                      prs.slide_width - Inches(1.4), Inches(1.8))
        tf = st.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.line_spacing = 1.35
        run(p, sub, size=22, color=sub_color)


# ───────────────────────── Slides ─────────────────────────

def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # Plain white background — no full-bleed colour.
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.5),
                              prs.slide_width, Inches(0.04))
    fill(band, TEAL)

    eyebrow = s.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(11), Inches(0.4))
    p = eyebrow.text_frame.paragraphs[0]
    run(p, "PREDICTIVE REVENUE INTELLIGENCE — CREDIT CARDS MVP",
        size=12, bold=True, color=TEAL)

    title = s.shapes.add_textbox(Inches(0.7), Inches(2.1), Inches(11), Inches(1.4))
    tf = title.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    run(p, "Predictive RPC Estimator", size=44, bold=True, color=NAVY)
    p2 = tf.add_paragraph()
    run(p2, "Real-time revenue-per-click forecasting for Credit Cards",
        size=20, color=SLATE)

    sub = s.shapes.add_textbox(Inches(0.7), Inches(4.8), Inches(11), Inches(2))
    tf = sub.text_frame; tf.word_wrap = True
    for label, val in [
        ("Solution demo", "Live staging environment"),
        ("Audience", "Credit Cards leadership + technical stakeholders"),
        ("Date", "2026-05-19"),
    ]:
        p = tf.add_paragraph(); p.line_spacing = 1.3
        run(p, f"{label}    ", size=12, bold=True, color=TEAL)
        run(p, val, size=14, color=SLATE)
    p = tf.add_paragraph(); p.line_spacing = 1.5
    run(p, " ", size=8)
    p = tf.add_paragraph()
    run(p, "Live dashboard:  ", size=11, bold=True, color=TEAL)
    run(p, DEMO_URL, size=11, color=MUTED)


def slide_opportunity(prs):
    """Sharp, hero-statement framing. White background, brand-blue accents.
    The villain (flat tCPA) named once, the cost named once."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    hero_statement(
        s, prs,
        big="Today every Credit Cards click is priced the same.",
        sub=(
            "You overpay for the minnow. You underbid the whale and a rival takes them. "
            "You see the bill 90 days later — when the ledger reconciles, and the budget is gone."
        ),
        eyebrow="The opportunity",
    )

    # Three short kickers along the bottom, on white.
    bottom = s.shapes.add_textbox(Inches(0.7), Inches(6.10),
                                  prs.slide_width - Inches(1.4), Inches(0.7))
    tf = bottom.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.line_spacing = 1.30
    for i, (eyebrow, text) in enumerate([
        ("WHY CC", "Near-term commercial pressure on the channel."),
        ("WHY NOW", "Sales coverage rising 50% → 80% by end-June."),
        ("WHY SAFE", "Bid-optimisation only — model output never touches customer terms."),
    ]):
        if i > 0:
            run(p, "      ", size=11, color=MUTED)
        run(p, eyebrow + "  ", size=10, bold=True, color=TEAL)
        run(p, text, size=11, color=SLATE)

    footer(s, prs, 2)


def slide_what_we_built(prs):
    """Four bold claims as the emotional pitch; the technical proof is
    a single small-print line under each, not a parallel column of bullets."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, prs, "We've built it. It's live.",
              eyebrow="What you can see — and what's running underneath")

    claims = [
        ("A bidder with a crystal ball",
         "Every Credit Cards click is priced in under a second, on its actual likely revenue.",
         "scoring-api (Rust on Cloud Run)  ·  Vertex AI XGBoost endpoint  ·  p95 < 1 s on staging."),
        ("Every prediction explainable",
         "Plain-English attribution chart on every score — auditable for finance and compliance.",
         "Native SHAP via Vertex explanationSpec  ·  bid-optimisation only, never customer terms."),
        ("Behavioural intent, day one",
         "The user's session — calculator, guides, compare-N — drives the bid in real time.",
         "Phoebe / GA4 nightly rollup  ·  Vertex Feature Store at serving time."),
        ("A safety net you can demo",
         "Bounds, circuit breaker, kill switch, anomaly window — the failure modes have already been thought through.",
         "Six Cloud Monitoring alerts  ·  six runbooks  ·  rollback is a Vertex traffic-split."),
    ]

    # 2×2 grid of claim cards
    cw = (prs.slide_width - Inches(1.2) - Inches(0.25)) / 2
    ch = Inches(2.55)
    left0 = Inches(0.55); top0 = Inches(1.20); gap = Inches(0.20)

    for i, (head, body, proof) in enumerate(claims):
        row, col = i // 2, i % 2
        x = left0 + col * (cw + gap)
        y = top0 + row * (ch + gap)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch)
        card.adjustments[0] = 0.04
        fill(card, WHITE)
        card.line.color.rgb = RULE; card.line.width = Pt(0.75)
        # Left accent stripe in our spine palette
        stripe_color = [TEAL, ACCENT, GREEN, NAVY][i]
        stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), ch)
        fill(stripe, stripe_color)
        # Number badge in the corner
        num = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 x + cw - Inches(0.55), y + Inches(0.20),
                                 Inches(0.35), Inches(0.35))
        fill(num, stripe_color)
        tf = num.text_frame
        tf.margin_left = Pt(0); tf.margin_right = Pt(0)
        tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run(p, str(i + 1), size=14, bold=True, color=WHITE)

        tb = s.shapes.add_textbox(x + Inches(0.30), y + Inches(0.18),
                                  cw - Inches(0.85), ch - Inches(0.36))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.line_spacing = 1.10
        run(p, head, size=20, bold=True, color=NAVY)
        p2 = tf.add_paragraph(); p2.line_spacing = 1.30
        p2.space_before = Pt(6)
        run(p2, body, size=12, color=SLATE)
        p3 = tf.add_paragraph(); p3.line_spacing = 1.30
        p3.space_before = Pt(6)
        run(p3, "Underneath:  ", size=9, bold=True, color=stripe_color)
        run(p3, proof, size=9, color=MUTED)

    # Bottom URL strip
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.55), Inches(6.55),
                             prs.slide_width - Inches(1.1), Inches(0.45))
    fill(bar, NAVY)
    tf = bar.text_frame
    tf.margin_left = Inches(0.18); tf.margin_top = Inches(0.10)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run(p, "Live on staging  ·  europe-west2  ·  ", size=10, color=LIGHT)
    run(p, DEMO_URL, size=10, bold=True, color=ACCENT)
    footer(s, prs, 3)


def slide_architecture(prs):
    """Native-pptx 4-lane diagram tailored for Credit Cards. Replaces
    the V1 embedded PNG which doesn't show Phoebe / GA4, the per-segment
    drift path, or the breach scanner.
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, prs, "End-to-end architecture",
              eyebrow="Credit Cards stack on Google Cloud")

    # 4 lanes, left to right
    lanes = [
        ("Sources", NAVY, [
            ("CM360 click stream", "Pub/Sub → BigQuery"),
            ("Sales ledger", "BigQuery, multi-stage events"),
            ("GA4 events (Phoebe)", "Daily BigQuery export"),
        ]),
        ("Real-time hot path", TEAL, [
            ("scoring-api (Rust, Cloud Run)", "p95 < 1s, /v1/score + /v1/explain"),
            ("Safety net", "Bounds, breaker, kill switch, anomaly window"),
            ("Vertex AI endpoint", "XGBoost + native SHAP"),
            ("Feature store", "Phoebe lookup at scoring time"),
        ]),
        ("Reconciliation + ML", ACCENT, [
            ("BigQuery views", "predictions_vs_revenue, coverage_audit"),
            ("Dataform", "residuals_by_segment, drift_breaches"),
            ("ml-pipeline", "Vertex Pipelines retrain on tag"),
            ("breach_emitter", "Daily scheduler → log → alerts"),
        ]),
        ("Surfaces + activation", GREEN, [
            ("Executive dashboard", "React + nginx on Cloud Run"),
            ("Activation API", "SA360 / SSGTM / OCI push"),
            ("Cloud Monitoring", "6 alert policies → on-call"),
        ]),
    ]

    n = len(lanes)
    margin = Inches(0.4)
    avail_w = prs.slide_width - Inches(0.8)
    lane_w = (avail_w - margin * (n - 1)) / n
    top = Inches(1.20)
    lane_h = Inches(5.55)

    for i, (lane_title, accent, items) in enumerate(lanes):
        x = Inches(0.4) + i * (lane_w + margin)

        # Lane card
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, lane_w, lane_h)
        card.adjustments[0] = 0.04
        fill(card, WHITE)
        card.line.color.rgb = RULE; card.line.width = Pt(0.75)

        # Header stripe
        stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top, lane_w, Inches(0.5))
        fill(stripe, accent)
        hb = s.shapes.add_textbox(x + Inches(0.15), top + Inches(0.06),
                                  lane_w - Inches(0.3), Inches(0.4))
        p = hb.text_frame.paragraphs[0]
        run(p, lane_title.upper(), size=11, bold=True, color=WHITE)

        # Items
        item_top = top + Inches(0.65)
        item_h = Inches(1.1)
        item_gap = Inches(0.12)
        for j, (label, sub) in enumerate(items):
            iy = item_top + j * (item_h + item_gap)
            ic = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x + Inches(0.15), iy,
                                   lane_w - Inches(0.3), item_h)
            ic.adjustments[0] = 0.10
            fill(ic, LIGHT)
            ic.line.color.rgb = RULE; ic.line.width = Pt(0.5)
            tf = ic.text_frame
            tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.10)
            tf.margin_top = Inches(0.12); tf.margin_bottom = Inches(0.08)
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run(p, label, size=11, bold=True, color=NAVY)
            p2 = tf.add_paragraph(); p2.line_spacing = 1.15
            run(p2, sub, size=9, color=SLATE)

        # Arrow to next lane
        if i < n - 1:
            arrow_x1 = x + lane_w + Pt(2)
            arrow_x2 = x + lane_w + margin - Pt(2)
            arrow_y = top + lane_h / 2
            ln = s.shapes.add_connector(2, arrow_x1, arrow_y, arrow_x2, arrow_y)
            ln.line.color.rgb = MUTED; ln.line.width = Pt(2)

    footer(s, prs, 4)


def slide_demo_walkthrough(prs):
    """Seven beats. Each is a number + headline + one short subline.
    No paragraph prose — the demo itself does the talking."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, prs, "What you'll see, in seven beats",
              eyebrow="Seven minutes, one URL")

    beats = [
        ("01", "KPIs in pounds",
         "Five tiles. 90-day window. Product-type filter."),
        ("02", "Coverage, slice by slice",
         "Red bars are the slices that need backfill before retraining."),
        ("03", "Same click, two worlds",
         "Flat target-CPA vs value-based bid vs the £ difference, with a running counter."),
        ("04", "You drive the model",
         "Predict in CC language; the pipeline trace lights up live."),
        ("05", "The bouncer with a crystal ball",
         "Press Play — a user's intent builds up, the bid follows."),
        ("06", "Why the model said that",
         "Plain-English SHAP under every prediction."),
        ("07", "What we've designed against",
         "Three failure modes named; the verticals after CC named."),
    ]

    # 4-up × 2-down grid for 7 beats; the 8th cell stays empty for the URL strip
    cols = 4
    rows = 2
    grid_w = prs.slide_width - Inches(1.1)
    grid_h = Inches(4.50)
    cw = (grid_w - Inches(0.15) * (cols - 1)) / cols
    ch = (grid_h - Inches(0.20) * (rows - 1)) / rows
    left0 = Inches(0.55); top0 = Inches(1.35)

    for i, (num, head, sub) in enumerate(beats):
        col = i % cols; row = i // cols
        x = left0 + col * (cw + Inches(0.15))
        y = top0 + row * (ch + Inches(0.20))
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch)
        card.adjustments[0] = 0.06
        fill(card, WHITE)
        card.line.color.rgb = RULE; card.line.width = Pt(0.75)

        # Big number, top-left
        nb = s.shapes.add_textbox(x + Inches(0.18), y + Inches(0.10),
                                  Inches(1.0), Inches(0.7))
        p = nb.text_frame.paragraphs[0]
        run(p, num, size=28, bold=True, color=ACCENT)

        # Headline
        hb = s.shapes.add_textbox(x + Inches(0.18), y + Inches(0.95),
                                  cw - Inches(0.36), Inches(0.7))
        tf = hb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.line_spacing = 1.10
        run(p, head, size=15, bold=True, color=NAVY)

        # Subline
        sb = s.shapes.add_textbox(x + Inches(0.18), y + Inches(1.55),
                                  cw - Inches(0.36), ch - Inches(1.7))
        tf = sb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.line_spacing = 1.30
        run(p, sub, size=11, color=SLATE)

    # URL strip in the empty 8th slot
    x = left0 + 3 * (cw + Inches(0.15))
    y = top0 + 1 * (ch + Inches(0.20))
    url_card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch)
    url_card.adjustments[0] = 0.06
    fill(url_card, NAVY)
    tf = url_card.text_frame
    tf.margin_left = Inches(0.20); tf.margin_right = Inches(0.20)
    tf.margin_top = Inches(0.25)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run(p, "Live dashboard", size=11, bold=True, color=ACCENT)
    p2 = tf.add_paragraph(); p2.line_spacing = 1.30; p2.space_before = Pt(8)
    run(p2, DEMO_URL, size=10, color=WHITE, font="Consolas")
    p3 = tf.add_paragraph(); p3.line_spacing = 1.30; p3.space_before = Pt(10)
    run(p3, "One URL. The dashboard auto-calls /v1/score, /v1/explain and /coverage via nginx — nothing else to open.",
        size=10, color=LIGHT)

    footer(s, prs, 5)


def slide_tech_choices(prs):
    """Headline-led tech rationale. The conclusion is the bold line;
    the reasoning is one sentence beneath. Pre-empts 'are you sure
    you didn't just pick the default' questions."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, prs, "Why this stack",
              eyebrow="The headline answer to every CTO question")

    choices = [
        ("Rust on Cloud Run",
         "We need every millisecond Vertex doesn't take.",
         "GC jitter would eat the safety-net budget. Considered Go and Python and rejected both."),
        ("XGBoost on Vertex AI",
         "Tabular data + auditable predictions + a circuit breaker that has to compose.",
         "GBDT wins tabular. Native SHAP comes free. Transformers' stateful decoding breaks the breaker."),
        ("TypeScript + React + Vite",
         "One URL with realtime interactivity, no server-side rendering.",
         "Zod validates JSON at the boundary. The dashboard is a static-asset deploy."),
        ("Protobuf",
         "Wire stability across Rust ↔ Python; CI parity test guards drift.",
         "Considered JSON Schema and rejected — loses on-wire byte-stability."),
        ("Dataform on BigQuery",
         "Same primitives as dbt; one fewer tool for the client to install.",
         "Type-safe ref()s give us the DAG automatically. coverage_audit, drift_breaches, residuals all light up in order."),
        ("Workload Identity Federation",
         "No long-lived secrets in CI. No keys to rotate, leak, or audit.",
         "GitHub OIDC federates into GCP. Every IAM grant is in Terraform."),
    ]

    # 3-col × 2-row grid
    cols = 3; rows = 2
    grid_w = prs.slide_width - Inches(1.1)
    grid_h = Inches(5.40)
    cw = (grid_w - Inches(0.20) * (cols - 1)) / cols
    ch = (grid_h - Inches(0.20) * (rows - 1)) / rows
    left0 = Inches(0.55); top0 = Inches(1.25)

    for i, (head, why, detail) in enumerate(choices):
        col = i % cols; row = i // cols
        x = left0 + col * (cw + Inches(0.20))
        y = top0 + row * (ch + Inches(0.20))
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch)
        card.adjustments[0] = 0.06
        fill(card, WHITE)
        card.line.color.rgb = RULE; card.line.width = Pt(0.75)
        # Top accent
        stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cw, Inches(0.08))
        fill(stripe, TEAL if i % 2 == 0 else ACCENT)

        tb = s.shapes.add_textbox(x + Inches(0.20), y + Inches(0.25),
                                  cw - Inches(0.40), ch - Inches(0.40))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.line_spacing = 1.10
        run(p, head, size=16, bold=True, color=NAVY)
        p2 = tf.add_paragraph(); p2.space_before = Pt(8); p2.line_spacing = 1.30
        run(p2, why, size=12, bold=True, color=SLATE)
        p3 = tf.add_paragraph(); p3.space_before = Pt(8); p3.line_spacing = 1.30
        run(p3, detail, size=10, color=MUTED)

    footer(s, prs, 6)


def slide_engineering(prs):
    """Inspection-sticker layout. Four production-readiness blocks, each
    with a big checkmark and the two or three things that matter."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, prs, "Production-grade engineering",
              eyebrow="The boring slide that makes the system bulletproof")

    blocks = [
        ("Reliability", TEAL, [
            "Bounds → flat-tCPA fallback on every prediction",
            "Circuit breaker auto-recovers; kill switch is a single flag, no redeploy",
            "Per-call timeouts on model + warehouse — no request hangs",
        ]),
        ("Observability", ACCENT, [
            "Six Cloud Monitoring alerts: latency, 5xx, breaker, anomaly, per-segment drift, coverage drop",
            "Cloud Logging + Trace baked in",
            "SLO: 99.5% availability on /v1/score, 30-day rolling",
        ]),
        ("Security + compliance", GREEN, [
            "Per-service service accounts; least-privilege IAM in Terraform",
            "Workload Identity Federation — no long-lived keys in CI",
            "Bid-optimisation only — no customer-decisioning surface, model output never reaches customer terms",
            "GA4 PII boundary — hashed identifiers only, raw user IDs and PII params stripped at ingest",
        ]),
        ("Operability", NAVY, [
            "Tag-push CI → CD → smoke; rollback is a single Vertex traffic-split",
            "Six runbooks committed: breaker reset, model rollback, secret rotation, coverage audit, BQ schema migration, endpoint scale-down",
        ]),
    ]

    # 2x2 grid
    cw = (prs.slide_width - Inches(1.1) - Inches(0.25)) / 2
    ch = Inches(2.65)
    left0 = Inches(0.55); top0 = Inches(1.20); gap = Inches(0.25)

    for i, (label, color, items) in enumerate(blocks):
        row = i // 2; col = i % 2
        x = left0 + col * (cw + gap)
        y = top0 + row * (ch + gap)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch)
        card.adjustments[0] = 0.05
        fill(card, WHITE)
        card.line.color.rgb = RULE; card.line.width = Pt(0.75)

        # Big check badge
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    x + Inches(0.25), y + Inches(0.25),
                                    Inches(0.6), Inches(0.6))
        fill(badge, color)
        tf = badge.text_frame
        tf.margin_left = Pt(0); tf.margin_right = Pt(0)
        tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run(p, "✓", size=22, bold=True, color=WHITE)

        # Heading
        hb = s.shapes.add_textbox(x + Inches(1.00), y + Inches(0.30),
                                  cw - Inches(1.20), Inches(0.5))
        p = hb.text_frame.paragraphs[0]
        run(p, label, size=18, bold=True, color=NAVY)

        # Items
        ib = s.shapes.add_textbox(x + Inches(0.25), y + Inches(1.00),
                                  cw - Inches(0.50), ch - Inches(1.10))
        tf = ib.text_frame; tf.word_wrap = True
        for j, it in enumerate(items):
            pp = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            pp.line_spacing = 1.25; pp.space_after = Pt(3)
            run(pp, "•  ", size=11, bold=True, color=color)
            run(pp, it, size=10, color=SLATE)

    footer(s, prs, 7)


def slide_data_model(prs):
    """CC schema + Phoebe + reconciliation + drift in one slide."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, prs, "Data and model lifecycle",
              eyebrow="From your ledger to a live prediction")
    body_box(s, prs, [
        ("h", "Three feeds, one model"),
        ("b", "Click stream (CM360, CC schema): product_type, card_product_id, query_intent, affinity_score, prior_applicant, income_band_bucket, rpc_14d/60d."),
        ("b", "Sales ledger: multi-stage events (application_started → submitted → approved → activated → first_spend → chargeback) with revenue, margin_rate (profit-ready), card_product_id, currency."),
        ("b", "GA4 / Phoebe: per-cookie 30-day rollup — calculator_used, guides_read, cards_compared, session_engagement_s. Joined at training; looked up at serving."),
        ("h", "Label and reconciliation window"),
        ("b", "Sum-of-rewards label over a 90-day window — fits the Credit Cards consideration tail. Profit-ready: when the commission table lands the same view becomes margin-aware with no rewrite."),
        ("h", "Training and release"),
        ("b", "Vertex AI Pipelines run from the same monorepo — reproducible from a commit. Every version in the Vertex Model Registry."),
        ("b", "Canary 10% → 50% over 48h → 100%. Active-versions panel shows the share and rolling MAE side-by-side."),
        ("h", "Drift, four ways"),
        ("b", "Inputs (PSI per feature); outputs (residuals_daily); per-segment MAE (product × device × geo); coverage drop W-o-W per slice."),
        ("mute", "Today's model is trained on synthetic data mirroring the CC schema. First work item once OQ-11 (GA4 access) and the data contract land is retrain on real."),
    ], font_size=11)
    footer(s, prs, 8)


def slide_roadmap(prs):
    """Three phase cards, tightened so text fits inside 4.0" height."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, prs, "Delivery roadmap", eyebrow="From here to cutover")

    phases = [
        ("Phase 1", "Platform & schema",
         "Complete", GREEN, [
            "CC schema end-to-end",
            "Phoebe features wired",
            "Demo dashboard live",
            "Drift + coverage alerts",
        ]),
        ("Phase 2", "Real data",
         "Pending client", ACCENT, [
            "GA4 access (OQ-11)",
            "Sample data export",
            "v1 on 50% coverage",
            "Canary 10% → 100%",
        ]),
        ("Phase 3", "Cutover",
         "End-August 2026", TEAL, [
            "v2 on 80% coverage",
            "Rollback rehearsed",
            "Compliance sign-off",
            "Handover to on-call",
        ]),
    ]
    cw = Inches(4.05); ch = Inches(4.20); top = Inches(1.30)
    left0 = Inches(0.55); gap = Inches(0.25)
    for i, (phase, title, status, color, items) in enumerate(phases):
        x = left0 + i * (cw + gap)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, cw, ch)
        card.adjustments[0] = 0.05
        fill(card, WHITE)
        card.line.color.rgb = RULE; card.line.width = Pt(0.75)
        # Header stripe
        stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top, cw, Inches(0.55))
        fill(stripe, color)
        hb = s.shapes.add_textbox(x + Inches(0.2), top + Inches(0.07),
                                  cw - Inches(0.4), Inches(0.5))
        p = hb.text_frame.paragraphs[0]
        run(p, phase, size=11, bold=True, color=WHITE)
        p2 = hb.text_frame.add_paragraph()
        run(p2, title, size=15, bold=True, color=WHITE)
        # Status pill
        sb = s.shapes.add_textbox(x + Inches(0.2), top + Inches(0.72),
                                  cw - Inches(0.4), Inches(0.32))
        p = sb.text_frame.paragraphs[0]
        run(p, status.upper(), size=10, bold=True, color=color)
        # Items
        ib = s.shapes.add_textbox(x + Inches(0.25), top + Inches(1.15),
                                  cw - Inches(0.5), ch - Inches(1.30))
        tf = ib.text_frame; tf.word_wrap = True
        for j, it in enumerate(items):
            pp = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            pp.line_spacing = 1.25; pp.space_after = Pt(4)
            run(pp, "✓  " if status == "Complete" else "•  ",
                size=12, bold=True, color=color)
            run(pp, it, size=12, color=SLATE)

    note = s.shapes.add_textbox(Inches(0.55), Inches(5.70),
                                prs.slide_width - Inches(1.1), Inches(0.6))
    p = note.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run(p, "Timeline measured from data-sample arrival, not kickoff. "
        "Phase 2 starts the day GA4 access + a sample data export land.",
        size=10, italic=True, color=MUTED)
    footer(s, prs, 9)


def slide_commercials(prs):
    """Three asks, pyramid hierarchy. The first is critical-path; the
    others wait their turn."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, prs, "What we need from you",
              eyebrow="Three asks, in priority order")

    asks = [
        ("1", "Critical path", TEAL,
         "GA4 access to your Credit Cards property",
         "OQ-11 — read on the analytics_<property> BigQuery export. Every week this slides is a week off cutover."),
        ("2", "Working session", ACCENT,
         "Confirm the GA4 event taxonomy + click→cookie join key",
         "OQ-12 + OQ-13 — one 60-minute call validates our best-current-guess mappings against your real events."),
        ("3", "Sign-offs", GREEN,
         "GCP project decision + compliance sign-off on the bid-optimisation + GA4 PII boundaries",
         "OQ-1 — client project or namespaced env. The CD job is wired and gated. Two named contacts: one data owner + one engineering lead."),
    ]

    # Stacked rows
    top0 = Inches(1.25)
    row_h = Inches(1.30)
    gap = Inches(0.20)
    for i, (num, eyebrow, color, head, sub) in enumerate(asks):
        y = top0 + i * (row_h + gap)
        # Number block
        nb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.55), y, Inches(1.10), row_h)
        nb.adjustments[0] = 0.10
        fill(nb, color)
        tf = nb.text_frame
        tf.margin_left = Pt(0); tf.margin_right = Pt(0)
        tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run(p, num, size=44, bold=True, color=WHITE)

        # Body
        bb = s.shapes.add_textbox(Inches(1.85), y + Inches(0.10),
                                  prs.slide_width - Inches(2.45), row_h - Inches(0.20))
        tf = bb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.line_spacing = 1.10
        run(p, eyebrow.upper() + "    ", size=10, bold=True, color=color)
        run(p, head, size=18, bold=True, color=NAVY)
        p2 = tf.add_paragraph(); p2.space_before = Pt(8); p2.line_spacing = 1.30
        run(p2, sub, size=11, color=SLATE)

    # Bottom note
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.55), Inches(6.40),
                             prs.slide_width - Inches(1.1), Inches(0.55))
    fill(bar, NAVY)
    tf = bar.text_frame
    tf.margin_left = Inches(0.20); tf.margin_top = Inches(0.12)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run(p, "WHAT YOU GET    ", size=10, bold=True, color=ACCENT)
    run(p, "A production service on your data, your on-call team in the cockpit, signed handover by end-August 2026.",
        size=11, color=LIGHT)

    footer(s, prs, 10)


def slide_thanks(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # Plain white background — brand-blue band for visual punctuation.
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.6),
                              prs.slide_width, Inches(0.06))
    fill(band, TEAL)

    eb = s.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(12), Inches(0.4))
    p = eb.text_frame.paragraphs[0]
    run(p, "OVER TO YOU", size=13, bold=True, color=TEAL)

    t = s.shapes.add_textbox(Inches(0.7), Inches(2.1), Inches(12), Inches(1.4))
    p = t.text_frame.paragraphs[0]
    run(p, "Let's get the data flowing.", size=60, bold=True, color=NAVY)

    sub = s.shapes.add_textbox(Inches(0.7), Inches(4.0), Inches(12), Inches(1.5))
    tf = sub.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.line_spacing = 1.30
    run(p, "GA4 access this week. Event taxonomy next. ", size=20, color=SLATE)
    run(p, "Cutover end-August.", size=20, bold=True, color=TEAL)

    foot = s.shapes.add_textbox(Inches(0.7), Inches(6.7), Inches(12), Inches(0.4))
    p = foot.text_frame.paragraphs[0]
    run(p, f"Live dashboard  ·  {DEMO_URL}", size=11, color=MUTED)


def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_cover(prs)              # 1
    slide_opportunity(prs)        # 2
    slide_what_we_built(prs)      # 3
    slide_architecture(prs)       # 4
    slide_demo_walkthrough(prs)   # 5
    slide_tech_choices(prs)       # 6
    slide_engineering(prs)        # 7
    slide_data_model(prs)         # 8
    slide_roadmap(prs)            # 9
    slide_commercials(prs)        # 10
    slide_thanks(prs)             # 11

    prs.save(OUT)
    print(f"wrote {OUT}  ({OUT.stat().st_size:,} bytes, {len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
