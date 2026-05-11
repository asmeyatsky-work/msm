"""Generate docs/searce-build-overview.pptx — 12-slide deck describing
the Searce-internal Predictive RPC Estimator reference build, demo link
included on the cover and demo slides.

Idempotent: re-run to regenerate.
"""
from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

OUT = Path(__file__).parent / "searce-build-overview.pptx"
DEMO_URL = "https://scoring-api-staging-ifcjcfl7xa-nw.a.run.app"

NAVY    = RGBColor(0x0B, 0x1F, 0x3A)
TEAL    = RGBColor(0x0F, 0x6E, 0x7A)
SLATE   = RGBColor(0x33, 0x3F, 0x52)
LIGHT   = RGBColor(0xF4, 0xF6, 0xFA)
ACCENT  = RGBColor(0xD9, 0x73, 0x06)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
MUTED   = RGBColor(0x6B, 0x73, 0x80)


def set_text(run, text, size=18, bold=False, color=SLATE, font="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = color


def add_title_bar(slide, prs, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.5); tf.margin_top = Inches(0.18)
    p = tf.paragraphs[0]
    set_text(p.add_run(), title, size=24, bold=True, color=WHITE)


def add_footer(slide, prs, page_num):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.35),
        prs.slide_width, Inches(0.35),
    )
    bar.fill.solid(); bar.fill.fore_color.rgb = LIGHT
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.5); tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    set_text(
        p.add_run(),
        f"Predictive RPC Estimator — Searce reference build   |   "
        f"{DEMO_URL}   |   {page_num} / 12",
        size=10, color=MUTED,
    )


def add_body_text(slide, prs, lines, top=Inches(1.1), left=Inches(0.5),
                  width=None, height=None, font_size=14, line_spacing=1.15):
    width = width or prs.slide_width - Inches(1.0)
    height = height or prs.slide_height - Inches(1.5)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        if isinstance(line, tuple):
            kind, text = line
            if kind == "h":
                set_text(p.add_run(), text, size=font_size + 4, bold=True, color=TEAL)
                p.space_after = Pt(4)
            elif kind == "b":
                p.level = 0
                set_text(p.add_run(), "•  ", size=font_size, bold=True, color=ACCENT)
                set_text(p.add_run(), text, size=font_size, color=SLATE)
            elif kind == "sb":
                p.level = 1
                set_text(p.add_run(), "–  ", size=font_size - 1, color=MUTED)
                set_text(p.add_run(), text, size=font_size - 1, color=SLATE)
            elif kind == "p":
                set_text(p.add_run(), text, size=font_size, color=SLATE)
            elif kind == "mono":
                set_text(p.add_run(), text, size=font_size - 2, color=NAVY, font="Menlo")
        else:
            set_text(p.add_run(), line, size=font_size, color=SLATE)


def add_table(slide, prs, headers, rows, top=Inches(1.2), left=Inches(0.5),
              col_widths=None, font_size=12):
    n_cols = len(headers); n_rows = len(rows) + 1
    width = prs.slide_width - Inches(1.0)
    height = Inches(0.4) * n_rows
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = tbl_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        tf = cell.text_frame; tf.margin_left = Inches(0.08); tf.margin_top = Inches(0.04)
        p = tf.paragraphs[0]
        set_text(p.add_run(), h, size=font_size, bold=True, color=WHITE)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
            tf = cell.text_frame; tf.margin_left = Inches(0.08); tf.margin_top = Inches(0.04)
            tf.word_wrap = True
            p = tf.paragraphs[0]
            set_text(p.add_run(), str(val), size=font_size, color=SLATE)


def slide_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ========================================================================== #

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

# ---------- 1. Cover ---------- #
s = slide_blank(prs)
hero = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
hero.fill.solid(); hero.fill.fore_color.rgb = NAVY; hero.line.fill.background()
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.4), prs.slide_width, Inches(0.18))
accent.fill.solid(); accent.fill.fore_color.rgb = ACCENT; accent.line.fill.background()

box = s.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(12), Inches(2.6))
tf = box.text_frame; tf.word_wrap = True
set_text(tf.paragraphs[0].add_run(), "Predictive RPC Estimator", size=52, bold=True, color=WHITE)
p = tf.add_paragraph()
set_text(p.add_run(), "Searce reference build — end-to-end on Google Cloud", size=24, color=LIGHT)
p = tf.add_paragraph(); p.space_before = Pt(36)
set_text(p.add_run(), "Live demo:  ", size=18, color=LIGHT)
set_text(p.add_run(), DEMO_URL, size=18, bold=True, color=ACCENT, font="Menlo")
p = tf.add_paragraph()
set_text(p.add_run(), "Built 2026-04   |   GCP project msm-rpc / europe-west2   |   tag v0.1.8",
         size=14, color=LIGHT)
add_footer(s, prs, 1)

# ---------- 2. Executive summary ---------- #
s = slide_blank(prs); add_title_bar(s, prs, "Executive summary")
add_body_text(s, prs, [
    ("h", "What this is"),
    ("p", "A working, end-to-end implementation of the Predictive RPC Estimator architecture from the PRD, "
          "deployed as a single Searce-internal reference environment on GCP. Every component is real "
          "(no mocks): live Vertex AI XGBoost endpoint, Cloud Run services, BigQuery analytics, "
          "Pub/Sub event flow, Cloud Monitoring alerts."),
    ("h", "Why we built it"),
    ("p", "To prove the architecture works as a single unit before quoting client deployment. The SOW "
          "for the client environment is built on the empirical numbers and integration patterns "
          "captured here — not on theoretical estimates."),
    ("h", "Status (2026-04-28)"),
    ("b", "All four Cloud Run services live; scoring-api serves real model + real explanations."),
    ("b", "Vertex AI XGBoost model trained, registered with explanationSpec, deployed."),
    ("b", "Predictions land in BigQuery via Pub/Sub→BQ subscription (verified end-to-end)."),
    ("b", "Four Cloud Monitoring alert policies wired to a verified email channel."),
    ("b", "CI/CD via GitHub Actions + Workload Identity Federation; no long-lived keys."),
    ("b", "Empirical load profile captured for SOW prod sizing."),
], top=Inches(1.05))
add_footer(s, prs, 2)

# ---------- 3. Architecture overview ---------- #
s = slide_blank(prs); add_title_bar(s, prs, "End-to-end architecture")

# layered diagram via shapes
def layer(left_in, top_in, w_in, h_in, label, fill, font_color=WHITE, font_size=12):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(left_in), Inches(top_in), Inches(w_in), Inches(h_in))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill; sh.line.fill.background()
    tf = sh.text_frame; tf.margin_left = Inches(0.08); tf.margin_top = Inches(0.04)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    set_text(p.add_run(), label, size=font_size, bold=True, color=font_color)
    return sh

# row 1: edge / client
layer(0.6, 1.2, 12.1, 0.6, "Client (sGTM / Browser) → POST /v1/score, /v1/explain", NAVY)
# row 2: cloud run scoring-api
layer(0.6, 2.0, 12.1, 0.7, "Cloud Run · scoring-api  (Rust · Axum · WIF SA)", TEAL)
# row 3: dependencies fanout
layer(0.6,  2.9, 3.0, 1.6, "Vertex AI Endpoint\nXGBoost @ explanationSpec\n(:predict / :explain)", SLATE, font_size=11)
layer(3.8,  2.9, 3.0, 1.6, "Pub/Sub\nrpc-predictions\nrpc-audit", SLATE, font_size=11)
layer(7.0,  2.9, 3.0, 1.6, "Secret Manager\nrpc-runtime-config\nvertex-endpoint-url", SLATE, font_size=11)
layer(10.2, 2.9, 2.5, 1.6, "BigQuery\nsales_ledger\n(breaker fallback)", SLATE, font_size=11)
# row 4: downstream sinks
layer(0.6, 4.7, 5.0, 0.7, "BigQuery rpc_predictions_raw  →  reconciliation views", TEAL)
layer(5.8, 4.7, 3.5, 0.7, "breaker-automation Cloud Run", TEAL)
layer(9.5, 4.7, 3.2, 0.7, "activation Cloud Run", TEAL)
# row 5: ops layer
layer(0.6, 5.6, 8.0, 0.7, "Cloud Monitoring · 4 alert policies → email channel (verified)", ACCENT, font_size=12)
layer(8.8, 5.6, 3.9, 0.7, "GitHub Actions CD (WIF · tag-driven)", ACCENT, font_size=12)

box = s.shapes.add_textbox(Inches(0.6), Inches(6.45), Inches(12.1), Inches(0.6))
tf = box.text_frame
set_text(tf.paragraphs[0].add_run(),
         "Single GCP project (msm-rpc, europe-west2). No mocks; "
         "every arrow above is a live RPC. min_instances=1 — single replica per service for cost discipline.",
         size=11, color=MUTED)
add_footer(s, prs, 3)

# ---------- 4. Live components inventory ---------- #
s = slide_blank(prs); add_title_bar(s, prs, "Live components")
add_table(s, prs,
    headers=["Layer", "Resource", "Purpose"],
    rows=[
        ["Compute",  "scoring-api-staging (Cloud Run, Rust)",   "POST /v1/score, /v1/explain — request path"],
        ["Compute",  "reconciliation-staging (Cloud Run)",      "Backfill predicted vs actual revenue from sales_ledger"],
        ["Compute",  "activation-staging (Cloud Run)",          "Push approved scores to sGTM / activation channels"],
        ["Compute",  "breaker-automation-staging (Cloud Run)",  "Subscribes rpc-anomaly → flips kill switch in Secret Manager"],
        ["ML",       "Vertex AI Endpoint 4471390533746425856",  "Online predict/explain on XGBoost rpc-estimator"],
        ["ML",       "Vertex AI Model 6679238469622956032",     "Trained on rpc_training_rows; sampled-shapley path=10"],
        ["Data",     "BigQuery rpc_estimator_staging dataset",  "8 tables/views — sales_ledger, rpc_predictions, …"],
        ["Event",    "Pub/Sub rpc-predictions / -audit / -clicks / -anomaly", "Decoupled fan-in / fan-out for guardrails"],
        ["Config",   "Secret Manager: rpc-runtime-config, vertex-endpoint-url, ssgtm-api-key", "Kill switch + per-env config without redeploy"],
        ["Identity", "Workload Identity Federation pool github-staging",         "Keyless GH Actions → ci-deployer-staging SA"],
        ["Monitor",  "4 google_monitoring_alert_policy + email channel",         "p95, 5xx, breaker, BQ-sub backlog"],
    ],
    col_widths=[1.6, 4.6, 6.1], font_size=11)
add_footer(s, prs, 4)

# ---------- 5. Request path ---------- #
s = slide_blank(prs); add_title_bar(s, prs, "Request path — what happens on a /v1/score")
add_body_text(s, prs, [
    ("h", "Domain-first construction"),
    ("b", "Request body is parsed into ClickFeatures via try_new — invalid inputs reject before any I/O (PRD §4)."),
    ("h", "Guardrails evaluated in order"),
    ("b", "Kill switch — if rpc-runtime-config.kill = true, return tCPA fallback immediately."),
    ("b", "Canary sampler — hash(click_id) decides if click is in canary band; out-of-canary skips the model."),
    ("b", "Circuit breaker — if Open and within cool-off, skip to data-layer fallback."),
    ("h", "Concurrent external calls (timeout-bounded)"),
    ("b", "Vertex :predict (model_timeout_ms) and Vertex CLV (clv_timeout) run via tokio::join!"),
    ("b", "Optional Feature Store enrichment with its own short timeout — degrades silently."),
    ("h", "Post-call invariants"),
    ("b", "Anomaly window records every model output; breached() trips the breaker."),
    ("b", "PredictionBounds rejects out-of-band values in favor of tCPA (PRD §5)."),
    ("h", "Side effects (best-effort, surfaced as warn! on failure as of v0.1.8)"),
    ("b", "Pub/Sub publish to rpc-predictions → BQ subscription → rpc_predictions_raw."),
    ("b", "Pub/Sub publish to rpc-audit (append-only)."),
], top=Inches(1.05), font_size=13)
add_footer(s, prs, 5)

# ---------- 6. ML pipeline ---------- #
s = slide_blank(prs); add_title_bar(s, prs, "ML model lifecycle")
add_body_text(s, prs, [
    ("h", "Training (one-shot, ops/deploy_real_model.py)"),
    ("b", "Fetches msm-rpc.rpc_estimator_staging.rpc_training_rows view from BigQuery (5,000 rows, synthetic)."),
    ("b", "Trains XGBoost regressor (max_depth=6, n_estimators=400, learning_rate=0.05, hist tree method)."),
    ("b", "Saves model.bst to gs://msm-rpc-rpc-artifacts-staging/models/rpc-estimator/<ts>/."),
    ("h", "Vertex AI registration — explanationSpec is mandatory (ADR 0002)"),
    ("b", "aiplatform.Model.upload with sampled-shapley parameters; path_count=10."),
    ("b", "explanation_metadata.json declares 8 input features under one tensor with index_feature_mapping."),
    ("h", "Online deploy"),
    ("b", "e2-standard-2 machine, min_replica_count=1 (Searce: cost discipline; client SOW will spec higher)."),
    ("b", "Scoring-api reads vertex-endpoint-url-staging secret; same URL serves :predict and :explain."),
    ("h", "Inference output (live)"),
    ("mono", '{"click_id":"…","predicted_rpc":7.054,"source":"Model"}'),
    ("mono", '{"base_value":-4.115,"contributions":[["rpc_7d",8.55],["cerberus_score",3.11],…]}'),
    ("h", "Phase 2.3 (client-gated, not in this build)"),
    ("b", "Re-train @2 against real cm360_clicks once data contract signed and ingestion lands."),
], top=Inches(1.05), font_size=13)
add_footer(s, prs, 6)

# ---------- 7. Guardrails ---------- #
s = slide_blank(prs); add_title_bar(s, prs, "PRD §5 guardrails — all live")
add_table(s, prs,
    headers=["Guardrail", "Behaviour", "How to demo it"],
    rows=[
        ["Kill switch",
         "Manual or breaker-driven flag in rpc-runtime-config secret. Forces tCPA fallback for every score within seconds.",
         "Bump secret version → observe scores switch to source=KILL_SWITCH."],
        ["Anomaly window",
         "Sliding/cumulative null-or-OOB rate; breach trips circuit breaker.",
         "Send 50+ requests with extreme features → window breached → breaker opens."],
        ["Circuit breaker",
         "Closed → Open → HalfOpen state machine; cool-off in millis. Open state skips model entirely.",
         "Force OPEN via anomaly → see source=FALLBACK_DATA_LAYER for cool-off duration."],
        ["Prediction bounds",
         "PredictionBounds.try_new from runtime-config; OOB outputs rejected before publish.",
         "Tighten bounds to {0.5, 2.0} via secret → high-rpc predictions take FALLBACK_TCPA path."],
        ["Canary sampler",
         "Hash(click_id) gates entry — out-of-canary clicks skip the model deterministically.",
         "Drop canary_bp to 5000 (50%) → ~half of distinct click_ids return FALLBACK_TCPA."],
        ["Reject-by-default constructors",
         "ClickFeatures::try_new and Rpc::try_new enforce invariants; bad input never touches Vertex.",
         "POST a malformed body → 4xx before any external RPC fires."],
    ],
    col_widths=[2.4, 5.0, 4.9], font_size=10.5)
add_footer(s, prs, 7)

# ---------- 8. Observability + Security ---------- #
s = slide_blank(prs); add_title_bar(s, prs, "Observability & security posture")
add_body_text(s, prs, [
    ("h", "Cloud Monitoring — wired and verified"),
    ("b", "scoring-api p95 latency > threshold (5m). Threshold is var-driven (staging=1500ms, prod default=1500ms)."),
    ("b", "scoring-api 5xx rate > 1% (5m)."),
    ("b", "scoring-api anomaly breaker tripped (any trip in 5m)."),
    ("b", "rpc-predictions BQ subscription backlog > 10k undelivered messages."),
    ("b", "Channel: email (verified). Confirmed end-to-end during this build by a real fire on 2026-04-28."),
    ("h", "Application telemetry"),
    ("b", "Prometheus /metrics endpoint scraped by Cloud Monitoring Managed Service. http_requests_total, "
          "http_request_errors_total, http_request_duration_seconds summary by path/status."),
    ("b", "Structured logs (tracing crate); sink-error visibility hardened in v0.1.8 — failures are now warn!, not silent."),
    ("h", "Security & supply chain"),
    ("b", "Workload Identity Federation: GitHub Actions assumes ci-deployer-staging via OIDC. No long-lived keys anywhere."),
    ("b", "Per-service runtime SAs (scoring-api, activation, breaker-automation) with least-privilege bindings."),
    ("b", "All runtime IAM grants codified in Terraform (main.tf + wif.tf) — reproducible from a fresh project."),
    ("b", "Secret Manager for runtime config; lifecycle{ ignore_changes } so terraform never overwrites breaker writes."),
    ("b", "CI: cargo build/test/clippy/fmt + python lint + proto parity + load-test smoke + e2e + SBOM (8 required checks)."),
], top=Inches(1.05), font_size=13)
add_footer(s, prs, 8)

# ---------- 9. Performance — measured ---------- #
s = slide_blank(prs); add_title_bar(s, prs, "Measured performance — load profile (staging)")
add_body_text(s, prs, [
    ("p", "30-second oha runs against the live staging URL, concurrency=10. Vertex XGBoost on e2-standard-2, "
          "scoring-api min=1. These are SOW inputs for client environment sizing."),
], top=Inches(1.0), font_size=12)
add_table(s, prs,
    headers=["Endpoint", "p50", "p95", "p99", "Errors", "Notes"],
    rows=[
        ["/v1/score",   "560-740 ms", "760-920 ms", "1.0-1.3 s", "3-10% (autoscale gap on min=1)",
         "Vertex round-trip dominates; XGBoost compute itself is ~few ms"],
        ["/v1/explain", "1.3-1.6 s",  "1.8-2.9 s",  "2.1-3.3 s", "<2% (post timeout bump v0.1.7)",
         "Sampled-shapley path=10 adds ~1s; intrinsic to algorithm"],
    ],
    top=Inches(1.65), col_widths=[1.6, 1.5, 1.5, 1.4, 3.0, 3.3], font_size=11)
add_body_text(s, prs, [
    ("h", "Implications for the client SOW (not for this Searce env)"),
    ("b", "Vertex machine size is the primary lever — n2-standard-4 typically halves p50."),
    ("b", "Cloud Run min_instances ≥ 2 in production closes the 5xx autoscale gap observed at min=1."),
    ("b", "Per-call timeouts must match the chosen Vertex sizing (currently model_timeout_ms=1500 staging-realistic)."),
    ("b", "If sub-second explain is required, drop sampled-shapley path_count or compute attributions async."),
], top=Inches(4.4), font_size=12)
add_footer(s, prs, 9)

# ---------- 10. Demo walkthrough ---------- #
s = slide_blank(prs); add_title_bar(s, prs, "Demo walkthrough")
add_body_text(s, prs, [
    ("h", "Demo URL"),
    ("mono", DEMO_URL),
    ("h", "1. Health probe"),
    ("mono", "curl https://scoring-api-staging-ifcjcfl7xa-nw.a.run.app/health   →   200 ok"),
    ("h", "2. Score a click — real Vertex prediction"),
    ("mono", 'curl -X POST -H "Content-Type: application/json" \\'),
    ("mono", '  $URL/v1/score -d \'{"click_id":"demo-1","correlation_id":"demo",'),
    ("mono", '  "device":"mobile","geo":"GB","hour_of_day":12,"query_intent":"high",'),
    ("mono", '  "ad_creative_id":"cr-1","cerberus_score":0.7,"rpc_7d":5.2,"rpc_14d":5.0,'),
    ("mono", '  "rpc_30d":4.8,"is_payday_week":true,"auction_pressure":0.6,'),
    ("mono", '  "landing_path":"/home","visits_prev_30d":1200}\''),
    ("p", "→ {\"predicted_rpc\":7.054,\"source\":\"Model\",…}"),
    ("h", "3. Explain — real per-feature attributions"),
    ("mono", "curl -X POST … $URL/v1/explain -d <same body>"),
    ("p", "→ {\"base_value\":-4.115,\"contributions\":[[\"rpc_7d\",8.55],[\"cerberus_score\",3.11],…]}"),
    ("h", "4. Verify the prediction landed in BigQuery"),
    ("mono", "bq query 'SELECT * FROM `msm-rpc.rpc_estimator_staging.rpc_predictions_raw`"),
    ("mono", "          WHERE publish_time > TIMESTAMP_SUB(CURRENT_TIMESTAMP(), INTERVAL 5 MINUTE)'"),
    ("h", "5. Inspect /metrics for live counters"),
    ("mono", "curl $URL/metrics  |  grep http_requests_total"),
], top=Inches(1.05), font_size=12)
add_footer(s, prs, 10)

# ---------- 11. What was built today vs preceding work ---------- #
s = slide_blank(prs); add_title_bar(s, prs, "Build history & most recent fixes")
add_body_text(s, prs, [
    ("h", "Foundations (prior sessions, through v0.1.5)"),
    ("b", "Domain + application + infrastructure crates with reject-by-default constructors and 95/85% coverage floors."),
    ("b", "Five Cloud Run services + Vertex endpoint + BQ dataset + Pub/Sub topology + WIF + monitoring policies."),
    ("b", "Synthetic data seed (5,000 clicks + 2,973 ledger rows) so the system exercises real Vertex without client data."),
    ("h", "This session (v0.1.6 → v0.1.8, 2026-04-28)"),
    ("b", "v0.1.6 — VertexExplain parser fixed for the indexed-array response shape that Vertex actually returns when "
          "explanation_metadata uses index_feature_mapping. Live /v1/explain went from empty contributions → real attributions."),
    ("b", "v0.1.7 — Bumped explain client + outer timeouts 1500 → 3000 ms. /v1/explain success rate 78.9% → 98.2%. "
          "Parameterized the p95 alert threshold; ran first empirical load profile."),
    ("b", "v0.1.8 — Granted scoring-api SA roles/pubsub.publisher on rpc-predictions and rpc-audit (was missing since launch). "
          "Replaced silent let _ = on sink calls with warn! — the bare let _ pattern had hidden the missing IAM since deploy."),
    ("h", "Operational outputs created today"),
    ("b", "Cloud Monitoring email notification channel created and verified — alerting now delivers."),
    ("b", "Load profile JSONs checked in under ops/perf/ for SOW input."),
    ("b", "Memory + project documents updated with empirical numbers, scope rules, and the fix history."),
], top=Inches(1.05), font_size=13)
add_footer(s, prs, 11)

# ---------- 12. From demo to client SOW ---------- #
s = slide_blank(prs); add_title_bar(s, prs, "From this demo to the client deployment")
add_body_text(s, prs, [
    ("h", "What this Searce build proves"),
    ("b", "The architecture works as a single unit on real GCP — every component, every guardrail, every contract."),
    ("b", "We have empirical p50/p95/p99 numbers, not estimates, to drive client sizing conversations."),
    ("b", "We have a reproducible Terraform/CI-CD path that can be re-pointed at the client project unchanged."),
    ("h", "Scope explicitly out of this Searce build"),
    ("b", "No prod-grade Vertex sizing on Searce; that is a client-environment decision driven by their traffic."),
    ("b", "No real client data — synthetic seed only. Phase 2.3 retrain is gated on data contract sign-off."),
    ("b", "No multi-region / DR — single europe-west2 deployment for cost; client SLOs may require more."),
    ("h", "What the SOW will cover for the client environment"),
    ("b", "Vertex sizing (n2-standard-4 baseline) + Cloud Run min_instances + per-call timeouts re-calibrated to that machine."),
    ("b", "Real ingestion: cm360_clicks Pub/Sub publisher identity + sales_ledger BigQuery transfer config."),
    ("b", "Phase 2.3: re-train against real data, register @2, traffic-split 10% → 100% with the breaker watching."),
    ("b", "Optional: SLOs + budget alerts, per-region replication, custom KMS keys, VPC-SC perimeter."),
    ("h", "Demo link"),
    ("mono", DEMO_URL),
    ("p", "Repo: https://github.com/asmeyatsky-work/msm   ·   Tag: v0.1.8   ·   GCP project: msm-rpc / europe-west2"),
], top=Inches(1.05), font_size=13)
add_footer(s, prs, 12)

prs.save(str(OUT))
print(f"wrote {OUT} ({OUT.stat().st_size} bytes)")
